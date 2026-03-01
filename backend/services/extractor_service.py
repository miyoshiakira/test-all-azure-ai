import io
from typing import List, Tuple
from dataclasses import dataclass


@dataclass
class Chunk:
    """チャンク（分割されたテキスト単位）"""
    text: str
    chunk_id: str  # 例: "page_1", "slide_3", "sheet_売上"
    chunk_type: str  # 例: "page", "slide", "sheet", "section"
    category: str = ""  # カテゴリ（仕事, 家族, 趣味 など）- 後でAIが設定


class TextExtractor:
    """各種ファイル形式からテキストを抽出するサービス"""

    # チャンクの設定（セマンティック検索最適化: 1チャンク = 1トピック）
    MAX_CHUNK_SIZE = 200      # チャンクの最大文字数（小さくして1トピックに）
    MIN_CHUNK_SIZE = 50       # チャンクの最小文字数
    OVERLAP_SIZE = 30         # チャンク間のオーバーラップ文字数

    @staticmethod
    def extract_chunks(file_content: bytes, file_name: str, content_type: str = "") -> Tuple[List[Chunk], str]:
        """
        ファイルからテキストをチャンク単位で抽出
        Returns: (chunks, file_type)
        """
        file_name_lower = file_name.lower()

        try:
            # PDF - ページごとに分割
            if file_name_lower.endswith('.pdf') or 'pdf' in content_type:
                return TextExtractor._extract_pdf_chunks(file_content), "PDF"

            # PowerPoint - スライドごとに分割
            elif file_name_lower.endswith(('.pptx', '.ppt')):
                return TextExtractor._extract_pptx_chunks(file_content), "PowerPoint"

            # Word - セクションごとに分割
            elif file_name_lower.endswith(('.docx', '.doc')):
                return TextExtractor._extract_docx_chunks(file_content), "Word"

            # Excel - シートごとに分割
            elif file_name_lower.endswith(('.xlsx', '.xls')):
                return TextExtractor._extract_xlsx_chunks(file_content), "Excel"

            # Text files - サイズで分割
            elif file_name_lower.endswith(('.txt', '.md', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.c', '.cpp', '.ts', '.tsx')):
                text = file_content.decode('utf-8', errors='ignore')
                return TextExtractor._split_text_to_chunks(text), "Text"

            # Try to decode as text
            elif 'text' in content_type:
                text = file_content.decode('utf-8', errors='ignore')
                return TextExtractor._split_text_to_chunks(text), "Text"

            # Unknown binary
            else:
                return [Chunk(f"[対応していないファイル形式: {file_name}]", "unknown", "unknown")], "Unknown"

        except Exception as e:
            return [Chunk(f"[テキスト抽出エラー: {str(e)}]", "error", "error")], "Error"

    @staticmethod
    def _extract_pdf_chunks(content: bytes) -> List[Chunk]:
        """PDFからページごとにチャンク抽出"""
        import fitz  # PyMuPDF

        chunks = []
        with fitz.open(stream=content, filetype="pdf") as doc:
            total_pages = len(doc)
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text().strip()
                if page_text:
                    # ページが長すぎる場合はさらに分割
                    if len(page_text) > TextExtractor.MAX_CHUNK_SIZE:
                        sub_chunks = TextExtractor._split_text_to_chunks(
                            page_text,
                            prefix=f"ページ {page_num}/{total_pages}"
                        )
                        for i, sub_chunk in enumerate(sub_chunks):
                            chunks.append(Chunk(
                                text=sub_chunk.text,
                                chunk_id=f"page_{page_num}_part_{i+1}",
                                chunk_type="page_section"
                            ))
                    else:
                        chunks.append(Chunk(
                            text=f"[ページ {page_num}/{total_pages}]\n\n{page_text}",
                            chunk_id=f"page_{page_num}",
                            chunk_type="page"
                        ))

        if not chunks:
            chunks.append(Chunk("[PDFからテキストを抽出できませんでした]", "error", "error"))

        return chunks

    @staticmethod
    def _extract_pptx_chunks(content: bytes) -> List[Chunk]:
        """PowerPointからスライドごとにチャンク抽出"""
        from pptx import Presentation

        chunks = []
        prs = Presentation(io.BytesIO(content))
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                slide_content = "\n".join(slide_texts)
                chunks.append(Chunk(
                    text=f"[スライド {slide_num}/{total_slides}]\n\n{slide_content}",
                    chunk_id=f"slide_{slide_num}",
                    chunk_type="slide"
                ))

        if not chunks:
            chunks.append(Chunk("[PowerPointからテキストを抽出できませんでした]", "error", "error"))

        return chunks

    @staticmethod
    def _extract_docx_chunks(content: bytes) -> List[Chunk]:
        """Wordからセクションごとにチャンク抽出"""
        from docx import Document

        doc = Document(io.BytesIO(content))

        # 段落を収集
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        # テーブルも収集
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_text.append(row_text)
            if table_text:
                paragraphs.append("\n".join(table_text))

        if not paragraphs:
            return [Chunk("[Wordからテキストを抽出できませんでした]", "error", "error")]

        # 段落をチャンクにまとめる
        chunks = []
        current_text = ""
        chunk_num = 1

        for para in paragraphs:
            if len(current_text) + len(para) > TextExtractor.MAX_CHUNK_SIZE:
                if current_text:
                    chunks.append(Chunk(
                        text=f"[セクション {chunk_num}]\n\n{current_text}",
                        chunk_id=f"section_{chunk_num}",
                        chunk_type="section"
                    ))
                    chunk_num += 1
                current_text = para
            else:
                current_text += "\n\n" + para if current_text else para

        # 残りを追加
        if current_text:
            chunks.append(Chunk(
                text=f"[セクション {chunk_num}]\n\n{current_text}",
                chunk_id=f"section_{chunk_num}",
                chunk_type="section"
            ))

        return chunks

    @staticmethod
    def _extract_xlsx_chunks(content: bytes) -> List[Chunk]:
        """Excelからシートごとにチャンク抽出"""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        chunks = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows_text = []

            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    rows_text.append(" | ".join(row_values))

            if rows_text:
                sheet_content = "\n".join(rows_text)

                # シートが長すぎる場合は分割
                if len(sheet_content) > TextExtractor.MAX_CHUNK_SIZE:
                    sub_chunks = TextExtractor._split_text_to_chunks(
                        sheet_content,
                        prefix=f"シート: {sheet_name}"
                    )
                    for i, sub_chunk in enumerate(sub_chunks):
                        chunks.append(Chunk(
                            text=sub_chunk.text,
                            chunk_id=f"sheet_{sheet_name}_part_{i+1}",
                            chunk_type="sheet_section"
                        ))
                else:
                    chunks.append(Chunk(
                        text=f"[シート: {sheet_name}]\n\n{sheet_content}",
                        chunk_id=f"sheet_{sheet_name}",
                        chunk_type="sheet"
                    ))

        wb.close()

        if not chunks:
            chunks.append(Chunk("[Excelからテキストを抽出できませんでした]", "error", "error"))

        return chunks

    @staticmethod
    def _split_into_sentences(text: str) -> List[str]:
        """テキストを文単位で分割"""
        import re
        # 日本語と英語の文末を考慮
        # 。！？.!? の後にスペースや改行がある場合、または文末の場合に分割
        sentences = re.split(r'(?<=[。！？.!?])\s*', text)
        return [s.strip() for s in sentences if s.strip()]

    @staticmethod
    def _split_by_topic(text: str) -> List[str]:
        """テキストをトピック（セクション/改行/見出し）単位で分割"""
        import re

        # 優先度1: 見出しパターンで分割（■、●、【】、#、数字. など）
        # 優先度2: 空行（2連続改行）で分割
        # 優先度3: 単一改行で分割

        # 見出しパターン
        heading_pattern = r'\n(?=(?:■|●|◆|▼|【|#|[0-9０-９]+[.．、)）]|\d+\.\s))'

        # まず見出しで分割
        sections = re.split(heading_pattern, text)

        result = []
        for section in sections:
            section = section.strip()
            if not section:
                continue

            # セクションが大きすぎる場合は空行で分割
            if len(section) > TextExtractor.MAX_CHUNK_SIZE:
                paragraphs = section.split('\n\n')
                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue
                    # まだ大きすぎる場合は単一改行で分割
                    if len(para) > TextExtractor.MAX_CHUNK_SIZE:
                        lines = para.split('\n')
                        for line in lines:
                            line = line.strip()
                            if line:
                                result.append(line)
                    else:
                        result.append(para)
            else:
                result.append(section)

        return result

    @staticmethod
    def _split_text_to_chunks(text: str, prefix: str = "") -> List[Chunk]:
        """テキストをトピック単位で分割（セマンティック検索最適化）"""

        # 短いテキストはそのまま返す
        if len(text) <= TextExtractor.MAX_CHUNK_SIZE:
            header = f"[{prefix}]\n\n" if prefix else ""
            return [Chunk(text=f"{header}{text}", chunk_id="chunk_1", chunk_type="text")]

        chunks = []
        chunk_num = 1

        # トピック単位で分割
        topics = TextExtractor._split_by_topic(text)

        current_chunk = ""
        overlap_text = ""

        for topic in topics:
            # トピックがまだ大きすぎる場合は文単位で分割
            if len(topic) > TextExtractor.MAX_CHUNK_SIZE:
                sentences = TextExtractor._split_into_sentences(topic)

                for sentence in sentences:
                    test_text = current_chunk + (" " if current_chunk else "") + sentence

                    if len(test_text) > TextExtractor.MAX_CHUNK_SIZE and current_chunk:
                        header = f"[{prefix} - {chunk_num}]\n\n" if prefix else ""
                        full_text = overlap_text + current_chunk if overlap_text else current_chunk
                        chunks.append(Chunk(
                            text=f"{header}{full_text}",
                            chunk_id=f"chunk_{chunk_num}",
                            chunk_type="topic"
                        ))

                        overlap_text = current_chunk[-TextExtractor.OVERLAP_SIZE:] if len(current_chunk) > TextExtractor.OVERLAP_SIZE else current_chunk
                        overlap_text = "..." + overlap_text + " "

                        chunk_num += 1
                        current_chunk = sentence
                    else:
                        current_chunk = test_text
            else:
                # トピックをそのままチャンクとして追加（理想的なケース）
                test_text = current_chunk + ("\n\n" if current_chunk else "") + topic

                if len(test_text) > TextExtractor.MAX_CHUNK_SIZE and current_chunk:
                    # 現在のチャンクを保存
                    header = f"[{prefix} - {chunk_num}]\n\n" if prefix else ""
                    full_text = overlap_text + current_chunk if overlap_text else current_chunk
                    chunks.append(Chunk(
                        text=f"{header}{full_text}",
                        chunk_id=f"chunk_{chunk_num}",
                        chunk_type="topic"
                    ))

                    overlap_text = current_chunk[-TextExtractor.OVERLAP_SIZE:] if len(current_chunk) > TextExtractor.OVERLAP_SIZE else current_chunk
                    overlap_text = "..." + overlap_text + " "

                    chunk_num += 1
                    current_chunk = topic
                else:
                    current_chunk = test_text

        # 残りのテキストを追加
        if current_chunk and len(current_chunk) >= TextExtractor.MIN_CHUNK_SIZE:
            header = f"[{prefix} - {chunk_num}]\n\n" if prefix else ""
            full_text = overlap_text + current_chunk if overlap_text else current_chunk
            chunks.append(Chunk(
                text=f"{header}{full_text}",
                chunk_id=f"chunk_{chunk_num}",
                chunk_type="topic"
            ))
        elif current_chunk and chunks:
            last_chunk = chunks[-1]
            chunks[-1] = Chunk(
                text=last_chunk.text + "\n\n" + current_chunk,
                chunk_id=last_chunk.chunk_id,
                chunk_type=last_chunk.chunk_type
            )
        elif current_chunk:
            header = f"[{prefix}]\n\n" if prefix else ""
            chunks.append(Chunk(
                text=f"{header}{current_chunk}",
                chunk_id="chunk_1",
                chunk_type="topic"
            ))

        return chunks

    @staticmethod
    def extract(file_content: bytes, file_name: str, content_type: str = "") -> Tuple[str, str]:
        """
        後方互換性のため: 全テキストを1つの文字列として返す
        """
        chunks, file_type = TextExtractor.extract_chunks(file_content, file_name, content_type)
        full_text = "\n\n---\n\n".join(chunk.text for chunk in chunks)
        return full_text, file_type
